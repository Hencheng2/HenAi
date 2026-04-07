# mydocs.py - Document creation utilities extracted from app.py

import os
import json
import pandas as pd
import openpyxl
from io import BytesIO
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import textwrap
from pathlib import Path


class DocumentCreator:
    """Utility class for creating various types of documents"""
    
    def __init__(self, output_dir="generated_docs"):
        """Initialize with output directory"""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.temp_dir = Path("temp_docs")
        self.temp_dir.mkdir(exist_ok=True)
    
    def create_word_document(self, content, filename):
        """Create a Word document from markdown-style content"""
        doc = Document()
        
        # Parse content into sections
        lines = content.split('\n')
        for line in lines:
            if line.strip():
                if line.startswith('###'):
                    doc.add_heading(line.replace('###', '').strip(), level=1)
                elif line.startswith('##'):
                    doc.add_heading(line.replace('##', '').strip(), level=2)
                elif line.startswith('#'):
                    doc.add_heading(line.replace('#', '').strip(), level=3)
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                else:
                    p = doc.add_paragraph(line)
        
        output_path = self.output_dir / filename
        doc.save(str(output_path))
        return output_path
    
    def create_text_file(self, content, filename):
        """Create a plain text file"""
        output_path = self.output_dir / filename
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return output_path
    
    def create_excel_file(self, content, filename):
        """Create an Excel file from CSV-like content"""
        # Parse content into rows (simple CSV-like parsing)
        rows = []
        lines = content.strip().split('\n')
        for line in lines:
            if ',' in line:
                rows.append(line.split(','))
            else:
                rows.append([line])
        
        df = pd.DataFrame(rows)
        output_path = self.output_dir / filename
        df.to_excel(str(output_path), index=False, header=False)
        return output_path
    
    def create_csv_file(self, content, filename):
        """Create a CSV file"""
        output_path = self.output_dir / filename
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return output_path
    
    def create_powerpoint(self, content, filename):
        """Create a PowerPoint presentation from slide content"""
        prs = Presentation()
        slides_content = content.split('---')
        
        for i, slide_content in enumerate(slides_content):
            lines = slide_content.strip().split('\n')
            if not lines:
                continue
            
            if i == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                
                if title and lines:
                    title.text = lines[0]
                if subtitle and len(lines) > 1:
                    subtitle.text = '\n'.join(lines[1:5])
            else:
                # Content slide
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                
                if title and lines:
                    title.text = lines[0]
                if content_shape and len(lines) > 1:
                    content_shape.text = '\n'.join(lines[1:])
        
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        return output_path
    
    def create_image_from_text(self, content, filename):
        """Create an image from text content"""
        # Calculate image size based on text length
        lines = textwrap.wrap(content, width=60)
        img_height = max(400, len(lines) * 30 + 100)
        img = Image.new('RGB', (800, img_height), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to load a font, fallback to default
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf", 20)
        except:
            font = ImageFont.load_default()
        
        y = 50
        for line in lines:
            draw.text((50, y), line, fill='black', font=font)
            y += 30
        
        output_path = self.output_dir / filename
        img.save(str(output_path))
        return output_path
    
    def create_pdf_from_content(self, content, filename):
        """Create a PDF from text content"""
        output_path = self.output_dir / filename
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        lines = content.split('\n')
        for line in lines:
            if line.strip():
                if line.startswith('#'):
                    story.append(Paragraph(line.lstrip('#').strip(), styles['Heading1']))
                else:
                    story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 12))
        
        doc.build(story)
        return output_path
    
    def create_document(self, content, doc_type, filename, template_id=None):
        """
        Main method to create a document based on type
        
        Args:
            content: The content to put in the document
            doc_type: Type of document ('word', 'txt', 'excel', 'csv', 'ppt', 'image', 'pdf')
            filename: Desired filename (without extension)
            template_id: Optional template ID for formatting
        
        Returns:
            Path to created document
        """
        # Ensure filename has proper extension
        ext_map = {
            'word': '.docx',
            'txt': '.txt',
            'pdf': '.pdf',
            'excel': '.xlsx',
            'csv': '.csv',
            'ppt': '.pptx',
            'image': '.png'
        }
        
        ext = ext_map.get(doc_type, '.txt')
        full_filename = f"{filename}{ext}"
        
        # Enhance content with template-specific formatting
        if template_id:
            content = self._apply_template_formatting(content, doc_type, template_id)
        
        # Create document based on type
        if doc_type == 'word':
            return self.create_word_document(content, full_filename)
        elif doc_type == 'txt':
            return self.create_text_file(content, full_filename)
        elif doc_type == 'excel':
            return self.create_excel_file(content, full_filename)
        elif doc_type == 'csv':
            return self.create_csv_file(content, full_filename)
        elif doc_type == 'ppt':
            return self.create_powerpoint(content, full_filename)
        elif doc_type == 'image':
            return self.create_image_from_text(content, full_filename)
        elif doc_type == 'pdf':
            return self.create_pdf_from_content(content, full_filename)
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
    
    def _apply_template_formatting(self, content, doc_type, template_id):
        """Apply template-specific formatting instructions to content"""
        if template_id == 'professional':
            if doc_type == 'word':
                content += """
                
                Format this as a professional document with:
                - Clear headings (use ### for main sections)
                - Bullet points where appropriate
                - Professional tone
                - Proper spacing between sections"""
            elif doc_type == 'ppt':
                content += """
                
                Format this as a presentation with:
                - Title slide (presentation title)
                - 3-5 content slides with headings and bullet points
                - Closing slide with summary or call to action
                Separate slides with ---"""
            elif doc_type == 'excel':
                content += """
                
                Format this as structured data with:
                - Column headers as first row
                - Each row as a data entry
                - Use consistent formatting"""
        
        return content
    
    def get_download_url(self, filename, base_url="/api/docs/download/"):
        """Get download URL for a file"""
        return f"{base_url}{filename}"


# Convenience functions for quick document creation

def create_word_document(content, filename, output_dir="generated_docs"):
    """Quickly create a Word document"""
    creator = DocumentCreator(output_dir)
    return creator.create_word_document(content, filename)


def create_text_file(content, filename, output_dir="generated_docs"):
    """Quickly create a text file"""
    creator = DocumentCreator(output_dir)
    return creator.create_text_file(content, filename)


def create_excel_file(content, filename, output_dir="generated_docs"):
    """Quickly create an Excel file"""
    creator = DocumentCreator(output_dir)
    return creator.create_excel_file(content, filename)


def create_powerpoint(content, filename, output_dir="generated_docs"):
    """Quickly create a PowerPoint presentation"""
    creator = DocumentCreator(output_dir)
    return creator.create_powerpoint(content, filename)


def create_image_from_text(content, filename, output_dir="generated_docs"):
    """Quickly create an image from text"""
    creator = DocumentCreator(output_dir)
    return creator.create_image_from_text(content, filename)


def create_pdf_from_content(content, filename, output_dir="generated_docs"):
    """Quickly create a PDF from text"""
    creator = DocumentCreator(output_dir)
    return creator.create_pdf_from_content(content, filename)


# Example usage
if __name__ == "__main__":
    # Test the document creator
    creator = DocumentCreator()
    
    # Create a sample Word document
    sample_content = """# Welcome to My Document
## Introduction
This is a sample document created with the DocumentCreator class.

## Features
- Easy document creation
- Multiple format support
- Professional formatting

### Conclusion
Thank you for using this utility!"""
    
    doc_path = creator.create_document(sample_content, 'word', 'sample_document')
    print(f"Created Word document: {doc_path}")
    
    # Create a sample text file
    text_path = creator.create_document("Hello, World!", 'txt', 'hello_world')
    print(f"Created text file: {text_path}")