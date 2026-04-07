"""
Binary File Processor for HenAi
Extracts metadata, text, and structured data from various binary file formats
- Multi-backend audio processing with fallbacks
- OCR using EasyOCR (no external dependencies)
- Comprehensive file type support
"""

import io
import os
import tempfile
from typing import Dict, Any, Optional, Tuple

# ============= TRY IMPORTS WITH FALLBACKS =============

# Image processing
try:
    from PIL import Image, ImageOps, ImageEnhance
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: PIL/Pillow not available. Install with: pip install Pillow")

try:
    import exifread
    EXIF_AVAILABLE = True
except ImportError:
    EXIF_AVAILABLE = False

# OCR (Pure Python, no external dependencies)
try:
    import easyocr
    EASYOCR_AVAILABLE = True
    _easyocr_reader = None
except ImportError:
    EASYOCR_AVAILABLE = False
    print("Warning: EasyOCR not available. Install with: pip install easyocr")

# Audio processing - multiple backends
try:
    from pydub import AudioSegment
    PYDUB_AVAILABLE = True
except ImportError:
    PYDUB_AVAILABLE = False

try:
    import speech_recognition as sr
    SPEECH_RECOGNITION_AVAILABLE = True
except ImportError:
    SPEECH_RECOGNITION_AVAILABLE = False

try:
    import mutagen
    MUTAGEN_AVAILABLE = True
except ImportError:
    MUTAGEN_AVAILABLE = False

try:
    import audioread
    AUDIOREAD_AVAILABLE = True
except ImportError:
    AUDIOREAD_AVAILABLE = False

try:
    import librosa
    LIBROSA_AVAILABLE = True
except ImportError:
    LIBROSA_AVAILABLE = False

# PDF processing
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

# Spreadsheet processing
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

# Document processing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Archive processing
try:
    import zipfile
    import tarfile
    ARCHIVE_AVAILABLE = True
except ImportError:
    ARCHIVE_AVAILABLE = False

# Encoding detection
try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

# Whisper for advanced transcription (optional)
try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False

# Video processing
try:
    import cv2
    import numpy as np
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False


def get_easyocr_reader():
    """Lazy initialization of EasyOCR reader"""
    global _easyocr_reader
    if _easyocr_reader is None and EASYOCR_AVAILABLE:
        try:
            # Use CPU only, English language
            _easyocr_reader = easyocr.Reader(['en'], gpu=False)
            print("EasyOCR initialized successfully")
        except Exception as e:
            print(f"Failed to initialize EasyOCR: {e}")
    return _easyocr_reader


class BinaryProcessor:
    """Main processor for all binary file types"""
    
    def __init__(self):
        self.initialize_handlers()
    
    def initialize_handlers(self):
        """Initialize all format-specific handlers"""
        self.handlers = {
            'image': self.process_image,
            'audio': self.process_audio,
            'video': self.process_video,
            'pdf': self.process_pdf,
            'spreadsheet': self.process_spreadsheet,
            'word': self.process_word_document,
            'presentation': self.process_presentation,
            'archive': self.process_archive,
            'database': self.process_database,
            'text': self.process_text_file,
        }
    
    def process_file(self, file_content: bytes, filename: str) -> str:
        """
        Main entry point - processes any file and returns formatted text for AI
        """
        file_ext = filename.split('.')[-1].lower() if '.' in filename else ''
        
        # Build output header
        output = f"\n\n--- FILE: {filename} ---\n"
        output += f"Size: {len(file_content)} bytes\n"
        
        # Route to appropriate handler based on extension
        if file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp', 'ico']:
            output += self.process_image(file_content, filename)
        elif file_ext in ['mp3', 'wav', 'ogg', 'flac', 'm4a', 'aac', 'wma', 'opus']:
            output += self.process_audio(file_content, filename)
        elif file_ext in ['mp4', 'avi', 'mov', 'mkv', 'webm', 'flv', 'wmv']:
            output += self.process_video(file_content, filename)
        elif file_ext == 'pdf':
            output += self.process_pdf(file_content, filename)
        elif file_ext in ['xlsx', 'xls', 'csv', 'xlsm', 'xlsb']:
            output += self.process_spreadsheet(file_content, filename)
        elif file_ext in ['docx', 'doc', 'odt']:
            output += self.process_word_document(file_content, filename)
        elif file_ext in ['pptx', 'ppt', 'odp']:
            output += self.process_presentation(file_content, filename)
        elif file_ext in ['zip', 'rar', '7z', 'tar', 'gz', 'bz2', 'xz']:
            output += self.process_archive(file_content, filename)
        elif file_ext in ['db', 'sqlite', 'sqlite3', 'db3']:
            output += self.process_database(file_content, filename)
        elif file_ext in ['txt', 'md', 'py', 'js', 'html', 'css', 'json', 'xml', 
                          'java', 'c', 'cpp', 'h', 'hpp', 'rb', 'php', 'go', 'rs', 
                          'swift', 'kt', 'ts', 'jsx', 'tsx', 'vue']:
            output += self.process_text_file(file_content, filename)
        else:
            # Try text extraction as fallback
            text_result = self.try_extract_text(file_content)
            if text_result:
                output += f"\n--- EXTRACTED TEXT ---\n{text_result}\n--- END TEXT ---\n"
            else:
                output += f"\n[Binary file: {filename}]\n"
                output += "No further extraction available for this file type.\n"
        
        output += "--- END FILE ---\n\n"
        return output
    
    def process_image(self, content: bytes, filename: str) -> str:
        """Extract image metadata and perform OCR using EasyOCR (no external dependencies)"""
        output = "\n--- IMAGE ANALYSIS ---\n"
        
        if not PIL_AVAILABLE:
            output += "❌ Image processing not available (Pillow not installed)\n"
            output += "Install with: pip install Pillow\n"
            output += "--- END IMAGE ANALYSIS ---\n"
            return output
        
        try:
            img = Image.open(io.BytesIO(content))
            output += f"📐 Dimensions: {img.width}x{img.height}\n"
            output += f"🎨 Format: {img.format}\n"
            output += f"🖼️ Mode: {img.mode}\n"
            
            # EXIF data
            if EXIF_AVAILABLE:
                try:
                    with io.BytesIO(content) as f:
                        tags = exifread.process_file(f)
                        if tags:
                            output += "\n📷 EXIF DATA:\n"
                            for tag, value in list(tags.items())[:10]:
                                output += f"  • {tag}: {value}\n"
                except:
                    pass
            
            # OCR for text in images using EasyOCR
            if EASYOCR_AVAILABLE:
                try:
                    reader = get_easyocr_reader()
                    if reader:
                        # Scale image if too large (improves OCR speed)
                        if img.width > 1500 or img.height > 1500:
                            img.thumbnail((1500, 1500))
                            output += f"\n📏 Image scaled for OCR\n"
                        
                        # Convert PIL image to numpy array
                        import numpy as np
                        img_array = np.array(img)
                        
                        # Run OCR
                        output += "\n🔍 OCR PROCESSING:\n"
                        results = reader.readtext(img_array)
                        
                        if results:
                            extracted_text = []
                            high_confidence_text = []
                            
                            for (bbox, text, confidence) in results:
                                if confidence > 0.5:
                                    high_confidence_text.append(text)
                                extracted_text.append(text)
                            
                            if high_confidence_text:
                                full_text = ' '.join(high_confidence_text)
                                output += f"✅ Extracted {len(full_text)} characters (high confidence)\n"
                                output += f"\n📝 EXTRACTED TEXT:\n{full_text.strip()}\n"
                            elif extracted_text:
                                full_text = ' '.join(extracted_text)
                                output += f"⚠️ Extracted {len(full_text)} characters (low confidence)\n"
                                output += f"\n📝 EXTRACTED TEXT:\n{full_text.strip()}\n"
                            else:
                                output += "❌ No readable text detected in image\n"
                        else:
                            output += "❌ No text detected in the image\n"
                except Exception as e:
                    output += f"\n⚠️ OCR processing error: {str(e)}\n"
                    output += "Make sure EasyOCR is installed: pip install easyocr\n"
            else:
                output += "\n⚠️ EasyOCR not installed. Install with: pip install easyocr\n"
                output += "This will enable text extraction from images without external dependencies.\n"
            
            output += "--- END IMAGE ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing image: {str(e)}\n"
        
        return output
    
    def extract_ocr_text(self, image_content: bytes, filename: str) -> str:
        """
        Extract only OCR text from an image without all the metadata
        """
        if not EASYOCR_AVAILABLE:
            return "[EasyOCR not installed. Install with: pip install easyocr]"
    
        try:
            from PIL import Image
            import numpy as np
        
            img = Image.open(io.BytesIO(image_content))
        
            # Scale image if too large
            if img.width > 1500 or img.height > 1500:
                img.thumbnail((1500, 1500))
        
            img_array = np.array(img)
            reader = get_easyocr_reader()
        
            if reader:
                results = reader.readtext(img_array)
                if results:
                    extracted_text = []
                    for (bbox, text, confidence) in results:
                        if confidence > 0.3:  # Lower threshold for more text
                            extracted_text.append(text)
                
                    if extracted_text:
                        return ' '.join(extracted_text)
        
            return ""
        except Exception as e:
            print(f"OCR extraction error: {e}")
            return f"[OCR error: {str(e)}]"
    
    def process_audio(self, content: bytes, filename: str) -> str:
        """
        Extract audio metadata and transcribe speech with multiple fallback methods
        Tries: 1. Mutagen (metadata) → 2. Audioread (info) → 3. Pydub (properties) → 4. Whisper (transcription) → 5. SpeechRecognition
        """
        output = "\n--- AUDIO ANALYSIS ---\n"
        output += f"🎵 File: {filename}\n"
        output += f"📦 Size: {len(content)} bytes\n"
        
        temp_file_path = None
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix='.' + filename.split('.')[-1], delete=False) as tmp:
                tmp.write(content)
                tmp.flush()
                temp_file_path = tmp.name
            
            # ============= METHOD 1: Mutagen (Best for metadata) =============
            if MUTAGEN_AVAILABLE:
                try:
                    audio_file = mutagen.File(temp_file_path)
                    if audio_file:
                        output += "\n📋 METADATA (Mutagen):\n"
                        
                        # Get info
                        if hasattr(audio_file, 'info'):
                            info = audio_file.info
                            if hasattr(info, 'length'):
                                minutes = int(info.length // 60)
                                seconds = int(info.length % 60)
                                output += f"  • Duration: {minutes}:{seconds:02d} ({info.length:.2f} seconds)\n"
                            if hasattr(info, 'bitrate'):
                                output += f"  • Bitrate: {info.bitrate} bps\n"
                            if hasattr(info, 'sample_rate'):
                                output += f"  • Sample Rate: {info.sample_rate} Hz\n"
                            if hasattr(info, 'channels'):
                                output += f"  • Channels: {info.channels}\n"
                        
                        # Get tags
                        if hasattr(audio_file, 'tags') and audio_file.tags:
                            output += "\n🏷️ TAGS:\n"
                            for key, value in list(audio_file.tags.items())[:15]:
                                output += f"  • {key}: {value}\n"
                except Exception as e:
                    output += f"\n⚠️ Mutagen metadata extraction failed: {str(e)}\n"
            
            # ============= METHOD 2: Audioread (Fallback for audio info) =============
            if AUDIOREAD_AVAILABLE and not (MUTAGEN_AVAILABLE and 'Duration' in output):
                try:
                    with audioread.audio_open(temp_file_path) as f:
                        output += "\n📊 AUDIO INFO (Audioread):\n"
                        duration = f.duration
                        minutes = int(duration // 60)
                        seconds = int(duration % 60)
                        output += f"  • Duration: {minutes}:{seconds:02d} ({duration:.2f} seconds)\n"
                        output += f"  • Sample Rate: {f.samplerate} Hz\n"
                        output += f"  • Channels: {f.channels}\n"
                        if hasattr(f, 'bitrate'):
                            output += f"  • Bitrate: {f.bitrate} bps\n"
                except Exception as e:
                    output += f"\n⚠️ Audioread info extraction failed: {str(e)}\n"
            
            # ============= METHOD 3: Pydub (For additional properties) =============
            if PYDUB_AVAILABLE:
                try:
                    audio = AudioSegment.from_file(temp_file_path)
                    duration = len(audio) / 1000
                    minutes = int(duration // 60)
                    seconds = int(duration % 60)
                    output += "\n🎚️ AUDIO PROPERTIES (Pydub):\n"
                    output += f"  • Duration: {minutes}:{seconds:02d} ({duration:.2f} seconds)\n"
                    output += f"  • Channels: {audio.channels}\n"
                    output += f"  • Frame Rate: {audio.frame_rate} Hz\n"
                    output += f"  • Sample Width: {audio.sample_width} bytes\n"
                    output += f"  • Max Amplitude: {audio.max}\n"
                    output += f"  • RMS: {audio.rms:.2f}\n"
                except Exception as e:
                    output += f"\n⚠️ Pydub processing failed: {str(e)}\n"
            
            # ============= METHOD 4: Whisper (Best for transcription - offline) =============
            if WHISPER_AVAILABLE:
                try:
                    output += "\n🎙️ WHISPER TRANSCRIPTION (Offline):\n"
                    output += "Loading Whisper model (first time may take a moment)...\n"
                    model = whisper.load_model("base")
                    result = model.transcribe(temp_file_path, language="en")
                    if result and result.get("text"):
                        transcript = result["text"].strip()
                        output += f"✅ Transcription complete!\n"
                        output += f"\n📝 TRANSCRIPT:\n{transcript}\n"
                    else:
                        output += "❌ No speech detected\n"
                except Exception as e:
                    output += f"⚠️ Whisper transcription failed: {str(e)}\n"
                    output += "Install Whisper: pip install openai-whisper torch\n"
            
            # ============= METHOD 5: SpeechRecognition (Fallback - online) =============
            elif SPEECH_RECOGNITION_AVAILABLE and not WHISPER_AVAILABLE:
                try:
                    # Try to convert to WAV for better compatibility
                    if PYDUB_AVAILABLE:
                        try:
                            audio = AudioSegment.from_file(temp_file_path)
                            wav_io = io.BytesIO()
                            audio.export(wav_io, format="wav")
                            wav_io.seek(0)
                            audio_source = wav_io
                        except:
                            audio_source = temp_file_path
                    else:
                        audio_source = temp_file_path
                    
                    recognizer = sr.Recognizer()
                    with sr.AudioFile(audio_source) as source:
                        output += "\n🎙️ SPEECH RECOGNITION (Google):\n"
                        recognizer.adjust_for_ambient_noise(source, duration=0.5)
                        audio_data = recognizer.record(source, duration=30)
                        
                        try:
                            transcript = recognizer.recognize_google(audio_data)
                            if transcript and transcript.strip():
                                output += f"✅ Transcription complete!\n"
                                output += f"\n📝 TRANSCRIPT:\n{transcript.strip()}\n"
                            else:
                                output += "❌ No speech detected\n"
                        except sr.UnknownValueError:
                            output += "❌ Could not understand audio\n"
                        except sr.RequestError as e:
                            output += f"⚠️ Google Speech Recognition error: {str(e)}\n"
                except Exception as e:
                    output += f"\n⚠️ Speech recognition failed: {str(e)}\n"
                    output += "Install SpeechRecognition: pip install SpeechRecognition\n"
            
            # ============= METHOD 6: Librosa (Scientific analysis) =============
            if LIBROSA_AVAILABLE:
                try:
                    import numpy as np
                    y, sr_lib = librosa.load(temp_file_path, sr=None, duration=30)
                    output += "\n🔬 AUDIO ANALYSIS (Librosa):\n"
                    output += f"  • RMS Energy: {np.mean(librosa.feature.rms(y=y)):.4f}\n"
                    output += f"  • Zero Crossing Rate: {np.mean(librosa.feature.zero_crossing_rate(y)):.4f}\n"
                    try:
                        tempo, _ = librosa.beat.beat_track(y=y, sr=sr_lib)
                        output += f"  • Estimated Tempo: {tempo:.2f} BPM\n"
                    except:
                        pass
                except Exception as e:
                    pass  # Silent fail for librosa as it's optional
            
            # Summary of what was successful
            output += "\n📊 PROCESSING SUMMARY:\n"
            success_count = 0
            if MUTAGEN_AVAILABLE and 'METADATA' in output:
                output += "  ✓ Metadata extracted (Mutagen)\n"
                success_count += 1
            if AUDIOREAD_AVAILABLE and 'Audioread' in output:
                output += "  ✓ Basic info extracted (Audioread)\n"
                success_count += 1
            if PYDUB_AVAILABLE and 'Pydub' in output:
                output += "  ✓ Audio properties analyzed (Pydub)\n"
                success_count += 1
            if WHISPER_AVAILABLE and 'TRANSCRIPT' in output:
                output += "  ✓ Speech transcribed (Whisper)\n"
                success_count += 1
            elif SPEECH_RECOGNITION_AVAILABLE and 'TRANSCRIPT' in output:
                output += "  ✓ Speech transcribed (Google)\n"
                success_count += 1
            
            if success_count == 0:
                output += "  ⚠️ Limited information available. Install additional packages:\n"
                output += "     • pip install mutagen audioread (for metadata)\n"
                output += "     • pip install openai-whisper torch (for transcription)\n"
                output += "     • pip install pydub (for audio properties)\n"
            
        except Exception as e:
            output += f"\n❌ Critical error processing audio file: {str(e)}\n"
        finally:
            # Clean up temp file
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
        
        output += "--- END AUDIO ANALYSIS ---\n"
        return output
    
    def process_video(self, content: bytes, filename: str) -> str:
        """Extract video metadata using multiple methods"""
        output = "\n--- VIDEO ANALYSIS ---\n"
        output += f"🎬 File: {filename}\n"
        output += f"📦 Size: {len(content)} bytes\n"
        
        temp_file_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.' + filename.split('.')[-1], delete=False) as tmp:
                tmp.write(content)
                tmp.flush()
                temp_file_path = tmp.name
            
            # Try OpenCV for video properties
            if CV2_AVAILABLE:
                try:
                    import numpy as np
                    cap = cv2.VideoCapture(temp_file_path)
                    if cap.isOpened():
                        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                        fps = cap.get(cv2.CAP_PROP_FPS)
                        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                        duration = frame_count / fps if fps > 0 else 0
                        
                        output += "\n🎥 VIDEO PROPERTIES (OpenCV):\n"
                        output += f"  • Resolution: {width}x{height}\n"
                        output += f"  • FPS: {fps:.2f}\n"
                        output += f"  • Frame Count: {frame_count}\n"
                        minutes = int(duration // 60)
                        seconds = int(duration % 60)
                        output += f"  • Duration: {minutes}:{seconds:02d} ({duration:.2f} seconds)\n"
                        cap.release()
                except Exception as e:
                    output += f"\n⚠️ OpenCV processing failed: {str(e)}\n"
            
            # Try moviepy if available
            try:
                from moviepy.editor import VideoFileClip
                clip = VideoFileClip(temp_file_path)
                output += "\n🎞️ VIDEO PROPERTIES (MoviePy):\n"
                output += f"  • Duration: {clip.duration:.2f} seconds\n"
                output += f"  • FPS: {clip.fps}\n"
                output += f"  • Size: {clip.size}\n"
                output += f"  • Has Audio: {clip.audio is not None}\n"
                clip.close()
            except ImportError:
                pass
            except Exception as e:
                output += f"\n⚠️ MoviePy processing failed: {str(e)}\n"
            
        except Exception as e:
            output += f"\n❌ Error processing video: {str(e)}\n"
        finally:
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
        
        output += "--- END VIDEO ANALYSIS ---\n"
        return output
    
    def process_pdf(self, content: bytes, filename: str) -> str:
        """Extract text, tables, and metadata from PDFs"""
        output = "\n--- PDF ANALYSIS ---\n"
        
        if not PDFPLUMBER_AVAILABLE:
            output += "❌ PDF processing not available (pdfplumber not installed)\n"
            output += "Install with: pip install pdfplumber\n"
            output += "--- END PDF ANALYSIS ---\n"
            return output
        
        try:
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                output += f"📄 Pages: {len(pdf.pages)}\n"
                
                # Extract metadata
                if pdf.metadata:
                    output += "\n📋 METADATA:\n"
                    for key, value in pdf.metadata.items():
                        if value:
                            output += f"  • {key}: {value}\n"
                
                # Extract text from all pages (limit to first 10)
                full_text = ""
                for i, page in enumerate(pdf.pages[:10]):
                    page_text = page.extract_text()
                    if page_text:
                        full_text += f"\n--- PAGE {i+1} ---\n{page_text}\n"
                
                if full_text:
                    # Limit total text to 10000 chars
                    if len(full_text) > 10000:
                        full_text = full_text[:10000] + "\n\n[Content truncated...]"
                    output += f"\n📝 TEXT CONTENT:\n{full_text}\n"
            
            output += "--- END PDF ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing PDF: {str(e)}\n"
        
        return output
    
    def process_spreadsheet(self, content: bytes, filename: str) -> str:
        """Extract data from Excel spreadsheets"""
        output = "\n--- SPREADSHEET ANALYSIS ---\n"
        
        if not PANDAS_AVAILABLE:
            output += "❌ Spreadsheet processing not available (pandas not installed)\n"
            output += "Install with: pip install pandas openpyxl\n"
            output += "--- END SPREADSHEET ANALYSIS ---\n"
            return output
        
        try:
            # Try pandas for comprehensive analysis
            df_dict = pd.read_excel(io.BytesIO(content), sheet_name=None)
            output += f"📊 Sheets: {', '.join(list(df_dict.keys()))}\n"
            
            for sheet_name, df in list(df_dict.items())[:3]:  # Limit to first 3 sheets
                output += f"\n📑 SHEET: {sheet_name}\n"
                output += f"  • Dimensions: {df.shape[0]} rows x {df.shape[1]} columns\n"
                output += f"  • Columns: {', '.join(df.columns.astype(str)[:15])}\n"
                
                # Show first 5 rows as sample
                sample = df.head(5).to_string()
                output += f"\n  Sample data:\n{sample}\n"
                
                # Basic statistics for numeric columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    output += f"\n  Numeric summary:\n"
                    output += df[numeric_cols].describe().to_string()
                    output += "\n"
            
            if len(df_dict) > 3:
                output += f"\n... and {len(df_dict) - 3} more sheets\n"
            
            output += "--- END SPREADSHEET ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing spreadsheet: {str(e)}\n"
        
        return output
    
    def process_word_document(self, content: bytes, filename: str) -> str:
        """Extract text from Word documents"""
        output = "\n--- WORD DOCUMENT ANALYSIS ---\n"
        
        if not DOCX_AVAILABLE:
            output += "❌ Word document processing not available (python-docx not installed)\n"
            output += "Install with: pip install python-docx\n"
            output += "--- END WORD DOCUMENT ANALYSIS ---\n"
            return output
        
        try:
            doc = Document(io.BytesIO(content))
            output += f"📝 Paragraphs: {len(doc.paragraphs)}\n"
            
            # Extract text from all paragraphs
            text = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            if text:
                if len(text) > 10000:
                    text = text[:10000] + "\n\n[Content truncated...]"
                output += f"\n📄 TEXT CONTENT:\n{text}\n"
            
            # Tables
            if doc.tables:
                output += f"\n📊 Tables found: {len(doc.tables)}\n"
            
            output += "--- END WORD DOCUMENT ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing Word document: {str(e)}\n"
        
        return output
    
    def process_presentation(self, content: bytes, filename: str) -> str:
        """Extract content from PowerPoint presentations"""
        output = "\n--- PRESENTATION ANALYSIS ---\n"
        
        if not PPTX_AVAILABLE:
            output += "❌ PowerPoint processing not available (python-pptx not installed)\n"
            output += "Install with: pip install python-pptx\n"
            output += "--- END PRESENTATION ANALYSIS ---\n"
            return output
        
        try:
            prs = Presentation(io.BytesIO(content))
            output += f"📽️ Slides: {len(prs.slides)}\n"
            
            slide_text = []
            for i, slide in enumerate(prs.slides[:10]):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                if slide_content:
                    slide_text.append(f"\n--- SLIDE {i+1} ---\n" + '\n'.join(slide_content))
            
            if slide_text:
                full_text = ''.join(slide_text)
                if len(full_text) > 10000:
                    full_text = full_text[:10000] + "\n\n[Content truncated...]"
                output += f"\n📝 TEXT CONTENT:\n{full_text}\n"
            
            output += "--- END PRESENTATION ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing presentation: {str(e)}\n"
        
        return output
    
    def process_archive(self, content: bytes, filename: str) -> str:
        """List archive contents"""
        output = "\n--- ARCHIVE ANALYSIS ---\n"
        
        try:
            file_ext = filename.split('.')[-1].lower()
            
            if file_ext == 'zip':
                with zipfile.ZipFile(io.BytesIO(content)) as zf:
                    files = zf.namelist()
                    output += f"📦 Total files: {len(files)}\n"
                    output += "\n📋 FILE LIST:\n"
                    for f in files[:50]:
                        info = zf.getinfo(f)
                        size = info.file_size
                        output += f"  • {f} ({size:,} bytes)\n"
                    if len(files) > 50:
                        output += f"  ... and {len(files) - 50} more files\n"
            elif file_ext in ['tar', 'gz', 'bz2']:
                with tarfile.open(fileobj=io.BytesIO(content), mode='r:*') as tf:
                    files = tf.getnames()
                    output += f"📦 Total files: {len(files)}\n"
                    output += "\n📋 FILE LIST:\n"
                    for f in files[:50]:
                        output += f"  • {f}\n"
                    if len(files) > 50:
                        output += f"  ... and {len(files) - 50} more files\n"
            else:
                output += f"Archive format {file_ext} - size: {len(content)} bytes\n"
                output += "For full archive support, install: pip install patool\n"
            
            output += "--- END ARCHIVE ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing archive: {str(e)}\n"
        
        return output
    
    def process_database(self, content: bytes, filename: str) -> str:
        """Analyze SQLite databases"""
        output = "\n--- DATABASE ANALYSIS ---\n"
        
        try:
            import sqlite3
            with tempfile.NamedTemporaryFile(suffix='.db', delete=False) as tmp:
                tmp.write(content)
                tmp.flush()
                tmp_path = tmp.name
            
            try:
                conn = sqlite3.connect(tmp_path)
                cursor = conn.cursor()
                
                # Get all tables
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = cursor.fetchall()
                
                output += f"🗄️ Tables: {len(tables)}\n\n"
                
                for table in tables[:20]:  # Limit to 20 tables
                    table_name = table[0]
                    cursor.execute(f"PRAGMA table_info({table_name})")
                    columns = cursor.fetchall()
                    
                    output += f"📋 TABLE: {table_name}\n"
                    output += f"  • Columns: {len(columns)}\n"
                    for col in columns[:15]:
                        output += f"    - {col[1]} ({col[2]})\n"
                    
                    # Get row count
                    cursor.execute(f"SELECT COUNT(*) FROM {table_name}")
                    row_count = cursor.fetchone()[0]
                    output += f"  • Rows: {row_count:,}\n"
                    
                    # Show sample data
                    if row_count > 0:
                        cursor.execute(f"SELECT * FROM {table_name} LIMIT 3")
                        sample = cursor.fetchall()
                        output += f"\n  Sample rows:\n"
                        for row in sample[:3]:
                            output += f"    {row}\n"
                    output += "\n"
                
                conn.close()
            finally:
                os.unlink(tmp_path)
            
            output += "--- END DATABASE ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing database: {str(e)}\n"
        
        return output
    
    def process_text_file(self, content: bytes, filename: str) -> str:
        """Enhanced text file processing with encoding detection"""
        output = "\n--- TEXT FILE ANALYSIS ---\n"
        
        try:
            # Detect encoding
            if CHARDET_AVAILABLE:
                detection = chardet.detect(content)
                encoding = detection.get('encoding', 'utf-8')
                confidence = detection.get('confidence', 0)
                output += f"🔤 Encoding: {encoding} (confidence: {confidence:.2%})\n"
            else:
                encoding = 'utf-8'
                output += "🔤 Encoding detection not available (install chardet)\n"
            
            # Decode content
            text = content.decode(encoding, errors='replace')
            lines = text.split('\n')
            output += f"📄 Lines: {len(lines):,}\n"
            output += f"📝 Characters: {len(text):,}\n"
            
            # Show first 100 lines as sample
            output += f"\n📖 SAMPLE CONTENT:\n"
            sample_lines = lines[:100]
            output += '\n'.join(sample_lines)
            if len(lines) > 100:
                output += f"\n... and {len(lines) - 100} more lines\n"
            
            output += "--- END TEXT FILE ANALYSIS ---\n"
            
        except Exception as e:
            output += f"❌ Error processing text file: {str(e)}\n"
        
        return output
    
    def try_extract_text(self, content: bytes) -> Optional[str]:
        """Attempt to extract text from unknown file types"""
        try:
            # Try to decode as UTF-8 first
            text = content.decode('utf-8', errors='replace')
            # Check if it looks like text (mostly printable)
            printable_chars = sum(1 for c in text if c.isprintable() or c in '\n\r\t')
            if printable_chars / len(text) > 0.7 and len(text) > 100:
                # Remove non-printable characters
                text = ''.join(char for char in text if char.isprintable() or char in '\n\r\t')
                return text[:5000]
            
            # Try with encoding detection
            if CHARDET_AVAILABLE:
                detection = chardet.detect(content)
                if detection['encoding']:
                    text = content.decode(detection['encoding'], errors='replace')
                    if len(text) > 100:
                        return text[:5000]
        except:
            pass
        
        return None