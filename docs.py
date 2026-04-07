"""
docs.py - Universal Document Processing Module
A comprehensive interface for processing, creating, editing, and converting
virtually all document types across different domains.
"""

import os
import json
import csv
import sqlite3
import zipfile
import tarfile
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Union, List, Dict, Any, Optional, BinaryIO
from datetime import datetime
import hashlib
import io

# Core data processing
import pandas as pd
import numpy as np

# Text and encoding
import chardet

# Word Processing
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PDF Processing
try:
    import pypdf
    import pdfplumber
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Spreadsheets
try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, Reference
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    import xlwt
    XLS_AVAILABLE = True
except ImportError:
    XLS_AVAILABLE = False

# Presentations
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PptxRGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Images - Using EasyOCR instead of Tesseract
try:
    from PIL import Image, ImageDraw, ImageFont, ImageFilter
    PILLOW_AVAILABLE = True
    # Use EasyOCR for text extraction
    try:
        import easyocr
        EASYOCR_AVAILABLE = True
        # Global reader instance (lazy initialization)
        _easyocr_reader = None
        print("✅ EasyOCR available for image text extraction")
    except ImportError:
        EASYOCR_AVAILABLE = False
        print("⚠️ EasyOCR not installed. Install with: pip install easyocr")
except ImportError:
    PILLOW_AVAILABLE = False
    EASYOCR_AVAILABLE = False
    print("⚠️ Pillow not installed. Install with: pip install Pillow")

# Markup and Web
try:
    from bs4 import BeautifulSoup
    import lxml
    import markdown
    BEAUTIFULSOUP_AVAILABLE = True
except ImportError:
    BEAUTIFULSOUP_AVAILABLE = False

# Data formats
try:
    import yaml
    import toml
    YAML_AVAILABLE = True
    TOML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False
    TOML_AVAILABLE = False

# Archives
try:
    import py7zr
    SEVENZ_AVAILABLE = True
except ImportError:
    SEVENZ_AVAILABLE = False

# Scientific
try:
    import h5py
    H5PY_AVAILABLE = True
except ImportError:
    H5PY_AVAILABLE = False

# Database
try:
    from sqlalchemy import create_engine, text
    SQLALCHEMY_AVAILABLE = True
except ImportError:
    SQLALCHEMY_AVAILABLE = False

# Audio/Video
try:
    from moviepy.editor import VideoFileClip, AudioFileClip, concatenate_videoclips
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False

try:
    from pydub import AudioSegment
    PYDUB_AVAILABLE = True
except ImportError:
    PYDUB_AVAILABLE = False

try:
    import mutagen
    from mutagen.mp3 import MP3
    from mutagen.flac import FLAC
    MUTAGEN_AVAILABLE = True
except ImportError:
    MUTAGEN_AVAILABLE = False


def get_easyocr_reader(lang='en'):
    """Lazy initialization of EasyOCR reader"""
    global _easyocr_reader
    if _easyocr_reader is None and EASYOCR_AVAILABLE:
        try:
            import easyocr
            # Initialize with English and common languages
            _easyocr_reader = easyocr.Reader([lang, 'en'], gpu=False)
            print("✅ EasyOCR reader initialized")
        except Exception as e:
            print(f"⚠️ Failed to initialize EasyOCR: {e}")
            return None
    return _easyocr_reader


class DocumentProcessor:
    """
    Universal Document Processor - Handles creation, reading, editing, conversion
    of virtually all document types.
    """
    
    def __init__(self, output_dir: str = "processed_docs"):
        """
        Initialize the document processor.
        
        Args:
            output_dir: Directory for output files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.temp_dir = Path(tempfile.mkdtemp())
        self.supported_formats = self._get_supported_formats()
        
        # Print available features
        print(f"📁 Document Processor initialized")
        print(f"   Output directory: {self.output_dir}")
        print(f"   Available: Word={DOCX_AVAILABLE}, PDF={PDF_AVAILABLE}, Excel={OPENPYXL_AVAILABLE}")
        print(f"   Available: PPT={PPTX_AVAILABLE}, Images={PILLOW_AVAILABLE}, OCR={EASYOCR_AVAILABLE}")
        
    def _get_supported_formats(self) -> Dict[str, List[str]]:
        """Return dictionary of supported formats by category."""
        return {
            "text": [".txt", ".text", ".asc", ".log"],
            "word": [".docx", ".doc", ".odt", ".rtf"],
            "pdf": [".pdf", ".xps"],
            "spreadsheet": [".xlsx", ".xls", ".csv", ".tsv", ".ods"],
            "presentation": [".pptx", ".ppt", ".odp", ".key"],
            "image": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp", ".svg"],
            "markup": [".html", ".htm", ".xml", ".md", ".rst"],
            "data": [".json", ".yaml", ".yml", ".toml", ".ini"],
            "archive": [".zip", ".tar", ".gz", ".7z", ".rar"],
            "ebook": [".epub", ".mobi", ".pdf"],
            "database": [".db", ".sqlite", ".sqlite3"],
            "audio": [".mp3", ".wav", ".flac", ".aac", ".ogg"],
            "video": [".mp4", ".avi", ".mov", ".mkv", ".webm"],
            "cad": [".stl", ".obj", ".dxf"],
            "scientific": [".h5", ".hdf5", ".nc", ".mat"],
            "business": [".qfx", ".ofx"],
            "calendar": [".ics", ".ical", ".vcf"],
            "security": [".pem", ".crt", ".key", ".p12"],
        }
    
    # ==================== TEXT DOCUMENTS ====================
    
    def create_text_file(self, content: str, output_path: str, encoding: str = 'utf-8') -> str:
        """Create a plain text file."""
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding=encoding) as f:
            f.write(content)
        
        return str(output_file)
    
    def read_text_file(self, file_path: str, detect_encoding: bool = True) -> str:
        """Read a text file with optional encoding detection."""
        path = Path(file_path)
        
        if detect_encoding:
            with open(path, 'rb') as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'
        else:
            encoding = 'utf-8'
        
        with open(path, 'r', encoding=encoding, errors='replace') as f:
            return f.read()
    
    def append_to_text_file(self, file_path: str, content: str) -> None:
        """Append content to an existing text file."""
        with open(file_path, 'a', encoding='utf-8') as f:
            f.write(content)
    
    # ==================== WORD DOCUMENTS ====================
    
    def create_word_document(self, output_path: str, content: List[Dict[str, Any]]) -> str:
        """
        Create a Word document.
        
        Args:
            output_path: Path for the output file
            content: List of dicts with 'type' ('paragraph', 'heading', 'table') and data
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is not installed")
        
        doc = Document()
        
        for item in content:
            if item['type'] == 'paragraph':
                p = doc.add_paragraph(item['text'])
                if 'style' in item:
                    p.style = item['style']
            elif item['type'] == 'heading':
                doc.add_heading(item['text'], level=item.get('level', 1))
            elif item['type'] == 'table':
                data = item['data']
                table = doc.add_table(rows=len(data), cols=len(data[0]))
                table.style = 'Table Grid'
                for i, row in enumerate(data):
                    for j, cell_value in enumerate(row):
                        table.cell(i, j).text = str(cell_value)
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_file))
        
        return str(output_file)
    
    def read_word_document(self, file_path: str) -> Dict[str, Any]:
        """Extract text and structure from a Word document."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is not installed")
        
        doc = Document(file_path)
        
        result = {
            'paragraphs': [],
            'tables': [],
            'metadata': {
                'core_properties': {},
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables)
            }
        }
        
        # Extract paragraphs
        for para in doc.paragraphs:
            result['paragraphs'].append(para.text)
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            result['tables'].append(table_data)
        
        return result
    
    # ==================== PDF DOCUMENTS ====================
    
    def create_pdf(self, output_path: str, content: List[Dict[str, Any]], 
                   page_size: str = 'letter') -> str:
        """Create a PDF document."""
        if not PDF_AVAILABLE:
            raise ImportError("ReportLab is not installed")
        
        page_sizes = {'letter': letter, 'A4': A4}
        size = page_sizes.get(page_size, letter)
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        doc = SimpleDocTemplate(str(output_file), pagesize=size)
        story = []
        styles = getSampleStyleSheet()
        
        for item in content:
            if item['type'] == 'paragraph':
                story.append(Paragraph(item['text'], styles['Normal']))
                story.append(Spacer(1, 12))
            elif item['type'] == 'heading':
                story.append(Paragraph(item['text'], styles['Heading1']))
                story.append(Spacer(1, 12))
            elif item['type'] == 'spacer':
                story.append(Spacer(1, item.get('height', 12)))
        
        doc.build(story)
        return str(output_file)
    
    def read_pdf(self, file_path: str, extract_tables: bool = True) -> Dict[str, Any]:
        """Extract text, tables, and metadata from a PDF."""
        result = {
            'pages': [],
            'metadata': {},
            'tables': []
        }
        
        # Extract text and metadata using PyPDF2
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            result['metadata'] = reader.metadata
            result['page_count'] = len(reader.pages)
            
            for i, page in enumerate(reader.pages):
                result['pages'].append({
                    'page_number': i + 1,
                    'text': page.extract_text()
                })
        
        # Extract tables using pdfplumber if available and requested
        if extract_tables and PDF_AVAILABLE:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for table in tables:
                            result['tables'].append({
                                'page': i + 1,
                                'data': table
                            })
            except:
                pass
        
        return result
    
    def merge_pdfs(self, pdf_paths: List[str], output_path: str) -> str:
        """Merge multiple PDF files into one."""
        if not PDF_AVAILABLE:
            raise ImportError("pypdf is not installed")
        
        merger = pypdf.PdfMerger()
        
        for pdf_path in pdf_paths:
            merger.append(pdf_path)
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        merger.write(str(output_file))
        merger.close()
        
        return str(output_file)
    
    def split_pdf(self, pdf_path: str, output_dir: str, pages_per_file: int = 1) -> List[str]:
        """Split a PDF into multiple files."""
        if not PDF_AVAILABLE:
            raise ImportError("pypdf is not installed")
        
        output_files = []
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        with open(pdf_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            total_pages = len(reader.pages)
            
            for start_page in range(0, total_pages, pages_per_file):
                writer = pypdf.PdfWriter()
                end_page = min(start_page + pages_per_file, total_pages)
                
                for page_num in range(start_page, end_page):
                    writer.add_page(reader.pages[page_num])
                
                output_file = output_path / f"part_{start_page // pages_per_file + 1}.pdf"
                with open(output_file, 'wb') as out_f:
                    writer.write(out_f)
                output_files.append(str(output_file))
        
        return output_files
    
    # ==================== SPREADSHEETS ====================
    
    def create_excel(self, output_path: str, sheets: Dict[str, List[List[Any]]]) -> str:
        """Create an Excel file with multiple sheets."""
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl is not installed")
        
        wb = openpyxl.Workbook()
        
        for i, (sheet_name, data) in enumerate(sheets.items()):
            if i == 0:
                ws = wb.active
                ws.title = sheet_name[:31]  # Excel sheet name limit
            else:
                ws = wb.create_sheet(title=sheet_name[:31])
            
            for row_idx, row in enumerate(data, 1):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(output_file))
        
        return str(output_file)
    
    def read_excel(self, file_path: str, sheet_name: Optional[str] = None) -> Dict[str, pd.DataFrame]:
        """Read Excel file using pandas."""
        if sheet_name:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            return {sheet_name: df}
        else:
            return pd.read_excel(file_path, sheet_name=None)
    
    def read_csv(self, file_path: str, **kwargs) -> pd.DataFrame:
        """Read CSV file with flexible options."""
        return pd.read_csv(file_path, **kwargs)
    
    def create_csv(self, data: List[Dict[str, Any]], output_path: str) -> str:
        """Create CSV file from list of dictionaries."""
        if not data:
            raise ValueError("No data provided")
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', newline='', encoding='utf-8') as f:
            writer = csv.DictWriter(f, fieldnames=data[0].keys())
            writer.writeheader()
            writer.writerows(data)
        
        return str(output_file)
    
    # ==================== PRESENTATIONS ====================
    
    def create_presentation(self, output_path: str, slides: List[Dict[str, Any]]) -> str:
        """Create a PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed")
        
        prs = Presentation()
        
        for slide_data in slides:
            slide_layout = prs.slide_layouts[slide_data.get('layout', 0)]
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            if 'title' in slide_data and slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
            
            # Add content
            if 'content' in slide_data and len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = slide_data['content']
            
            # Add images if specified
            if 'images' in slide_data:
                for img_data in slide_data['images']:
                    left = PptxInches(img_data.get('left', 1))
                    top = PptxInches(img_data.get('top', 2))
                    slide.shapes.add_picture(img_data['path'], left, top,
                                            width=PptxInches(img_data.get('width', 3)))
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_file))
        
        return str(output_file)
    
    def read_presentation(self, file_path: str) -> Dict[str, Any]:
        """Extract content from a PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed")
        
        prs = Presentation(file_path)
        
        result = {
            'slides': [],
            'slide_count': len(prs.slides)
        }
        
        for slide in prs.slides:
            slide_data = {
                'shapes': [],
                'notes': None,
                'text_content': ''
            }
            
            # Extract text from shapes
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
                    slide_data['shapes'].append({
                        'type': type(shape).__name__,
                        'text': shape.text
                    })
            
            slide_data['text_content'] = '\n'.join(text_parts)
            
            # Extract notes if available
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_data['notes'] = notes_slide.notes_text_frame.text
            
            result['slides'].append(slide_data)
        
        return result
    
    # ==================== IMAGES ====================
    
    def convert_image(self, input_path: str, output_path: str, format: str = 'PNG', 
                     quality: int = 95) -> str:
        """Convert image between formats with optional quality setting."""
        if not PILLOW_AVAILABLE:
            raise ImportError("Pillow is not installed")
        
        img = Image.open(input_path)
        
        # Convert RGBA to RGB for JPEG
        if format.upper() == 'JPEG' and img.mode == 'RGBA':
            rgb_img = Image.new('RGB', img.size, (255, 255, 255))
            rgb_img.paste(img, mask=img.split()[3])
            img = rgb_img
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        if format.upper() in ['JPEG', 'JPG']:
            img.save(str(output_file), format=format.upper(), quality=quality)
        else:
            img.save(str(output_file), format=format.upper())
        
        return str(output_file)
    
    def resize_image(self, input_path: str, output_path: str, 
                    width: int = None, height: int = None, 
                    maintain_aspect: bool = True) -> str:
        """Resize an image."""
        if not PILLOW_AVAILABLE:
            raise ImportError("Pillow is not installed")
        
        img = Image.open(input_path)
        
        if maintain_aspect and width and height:
            ratio = min(width / img.width, height / img.height)
            new_width = int(img.width * ratio)
            new_height = int(img.height * ratio)
        elif width:
            ratio = width / img.width
            new_width = width
            new_height = int(img.height * ratio)
        elif height:
            ratio = height / img.height
            new_height = height
            new_width = int(img.width * ratio)
        else:
            new_width, new_height = img.width, img.height
        
        resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        resized_img.save(str(output_file))
        
        return str(output_file)
    
    def extract_text_from_image(self, image_path: str, lang: str = 'en') -> str:
        """Extract text from image using EasyOCR."""
        if not EASYOCR_AVAILABLE:
            return "[OCR not available. Please install EasyOCR: pip install easyocr]"
        
        try:
            # Get EasyOCR reader
            reader = get_easyocr_reader(lang)
            if reader is None:
                return "[Failed to initialize EasyOCR]"
            
            # Read text from image
            result = reader.readtext(image_path, detail=0)  # detail=0 returns only text
            
            if result:
                # Join all detected text
                extracted_text = '\n'.join(result)
                return extracted_text
            else:
                return "[No text detected in image]"
                
        except Exception as e:
            return f"[OCR error: {str(e)}]"
    
    def create_image_from_text(self, text: str, output_path: str, 
                               width: int = 800, height: int = 600,
                               bg_color: str = 'white', text_color: str = 'black') -> str:
        """Create an image with text."""
        if not PILLOW_AVAILABLE:
            raise ImportError("Pillow is not installed")
        
        img = Image.new('RGB', (width, height), color=bg_color)
        draw = ImageDraw.Draw(img)
        
        # Try to load a default font, fallback to default if not found
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()
        
        # Wrap text
        margin = 50
        max_width = width - 2 * margin
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=font)
            text_width = bbox[2] - bbox[0]
            if text_width <= max_width:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        # Draw text
        y = margin
        for line in lines:
            draw.text((margin, y), line, fill=text_color, font=font)
            y += 30
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        img.save(str(output_file))
        
        return str(output_file)
    
    # ==================== MARKUP DOCUMENTS ====================
    
    def parse_html(self, html_content: str) -> BeautifulSoup:
        """Parse HTML content."""
        if not BEAUTIFULSOUP_AVAILABLE:
            raise ImportError("BeautifulSoup4 is not installed")
        
        return BeautifulSoup(html_content, 'lxml')
    
    def create_html(self, title: str, body_content: str, output_path: str) -> str:
        """Create an HTML document."""
        html_template = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        h1 {{ color: #333; }}
        .container {{ max-width: 800px; margin: auto; }}
    </style>
</head>
<body>
    <div class="container">
        {body_content}
    </div>
</body>
</html>"""
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(html_template)
        
        return str(output_file)
    
    def convert_markdown_to_html(self, markdown_content: str, output_path: str) -> str:
        """Convert Markdown to HTML."""
        if not BEAUTIFULSOUP_AVAILABLE:
            raise ImportError("markdown library is not installed")
        
        html_content = markdown.markdown(markdown_content)
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return str(output_file)
    
    # ==================== DATA FORMATS ====================
    
    def read_json(self, file_path: str) -> Dict[str, Any]:
        """Read JSON file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def write_json(self, data: Dict[str, Any], output_path: str, indent: int = 2) -> str:
        """Write JSON file."""
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=indent)
        
        return str(output_file)
    
    def read_yaml(self, file_path: str) -> Dict[str, Any]:
        """Read YAML file."""
        if not YAML_AVAILABLE:
            raise ImportError("PyYAML is not installed")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)
    
    def write_yaml(self, data: Dict[str, Any], output_path: str) -> str:
        """Write YAML file."""
        if not YAML_AVAILABLE:
            raise ImportError("PyYAML is not installed")
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, default_flow_style=False)
        
        return str(output_file)
    
    def read_toml(self, file_path: str) -> Dict[str, Any]:
        """Read TOML file."""
        if not TOML_AVAILABLE:
            raise ImportError("toml library is not installed")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            return toml.load(f)
    
    # ==================== ARCHIVES ====================
    
    def create_zip_archive(self, files: List[str], output_path: str) -> str:
        """Create a ZIP archive from files."""
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with zipfile.ZipFile(output_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for file in files:
                file_path = Path(file)
                if file_path.exists():
                    zipf.write(file_path, file_path.name)
        
        return str(output_file)
    
    def extract_archive(self, archive_path: str, extract_to: str = None) -> str:
        """Extract various archive formats."""
        extract_path = Path(extract_to) if extract_to else self.output_dir / "extracted"
        extract_path.mkdir(parents=True, exist_ok=True)
        
        archive_path = Path(archive_path)
        
        if archive_path.suffix == '.zip':
            with zipfile.ZipFile(archive_path, 'r') as zipf:
                zipf.extractall(extract_path)
        
        elif archive_path.suffix in ['.tar', '.gz', '.bz2']:
            with tarfile.open(archive_path, 'r:*') as tarf:
                tarf.extractall(extract_path)
        
        elif archive_path.suffix == '.7z' and SEVENZ_AVAILABLE:
            with py7zr.SevenZipFile(archive_path, mode='r') as szf:
                szf.extractall(extract_path)
        
        return str(extract_path)
    
    # ==================== DATABASES ====================
    
    def execute_sql_query(self, db_path: str, query: str) -> List[tuple]:
        """Execute SQL query on SQLite database."""
        conn = sqlite3.connect(db_path)
        try:
            cursor = conn.cursor()
            cursor.execute(query)
            results = cursor.fetchall()
            conn.commit()
            return results
        finally:
            conn.close()
    
    def create_sqlite_db(self, output_path: str, tables: Dict[str, List[Dict[str, Any]]]) -> str:
        """Create SQLite database with tables."""
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        conn = sqlite3.connect(str(output_file))
        cursor = conn.cursor()
        
        for table_name, rows in tables.items():
            if rows:
                columns = rows[0].keys()
                create_sql = f"CREATE TABLE IF NOT EXISTS {table_name} ({', '.join(columns)})"
                cursor.execute(create_sql)
                
                for row in rows:
                    placeholders = ', '.join(['?' for _ in columns])
                    insert_sql = f"INSERT INTO {table_name} VALUES ({placeholders})"
                    cursor.execute(insert_sql, list(row.values()))
        
        conn.commit()
        conn.close()
        
        return str(output_file)
    
    # ==================== AUDIO ====================
    
    def convert_audio(self, input_path: str, output_path: str, format: str = 'mp3') -> str:
        """Convert audio between formats."""
        if not PYDUB_AVAILABLE:
            raise ImportError("pydub is not installed")
        
        audio = AudioSegment.from_file(input_path)
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        audio.export(str(output_file), format=format)
        
        return str(output_file)
    
    def get_audio_metadata(self, audio_path: str) -> Dict[str, Any]:
        """Extract metadata from audio files."""
        if not MUTAGEN_AVAILABLE:
            raise ImportError("mutagen is not installed")
        
        audio_path = Path(audio_path)
        metadata = {}
        
        try:
            if audio_path.suffix == '.mp3':
                audio = MP3(audio_path)
            elif audio_path.suffix == '.flac':
                audio = FLAC(audio_path)
            else:
                audio = mutagen.File(audio_path)
            
            if audio:
                metadata['length'] = audio.info.length
                metadata['bitrate'] = getattr(audio.info, 'bitrate', None)
                metadata['sample_rate'] = getattr(audio.info, 'sample_rate', None)
                metadata['channels'] = getattr(audio.info, 'channels', None)
                
                # Tags
                if hasattr(audio, 'tags'):
                    metadata['tags'] = dict(audio.tags)
        
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata
    
    # ==================== VIDEO ====================
    
    def extract_video_info(self, video_path: str) -> Dict[str, Any]:
        """Extract information from video file."""
        if not MOVIEPY_AVAILABLE:
            # Try to get basic info using other methods
            try:
                import subprocess
                result = subprocess.run(['ffprobe', '-v', 'quiet', '-print_format', 'json', '-show_format', '-show_streams', video_path], 
                                      capture_output=True, text=True)
                if result.returncode == 0:
                    import json
                    data = json.loads(result.stdout)
                    info = {
                        'duration': float(data.get('format', {}).get('duration', 0)),
                        'size': int(data.get('format', {}).get('size', 0)),
                        'filename': Path(video_path).name
                    }
                    # Try to get video streams
                    for stream in data.get('streams', []):
                        if stream.get('codec_type') == 'video':
                            info['width'] = stream.get('width')
                            info['height'] = stream.get('height')
                            info['fps'] = stream.get('r_frame_rate')
                            break
                    return info
            except:
                pass
            return {
                'duration': 0,
                'size': 0,
                'filename': Path(video_path).name,
                'error': 'MoviePy not installed, video analysis limited'
            }
        
        try:
            clip = VideoFileClip(video_path)
            
            info = {
                'duration': clip.duration,
                'size': clip.size,
                'fps': clip.fps,
                'filename': Path(video_path).name
            }
            
            clip.close()
            return info
        except Exception as e:
            return {
                'filename': Path(video_path).name,
                'error': str(e)
            }
    
    def extract_audio_from_video(self, video_path: str, output_path: str) -> str:
        """Extract audio track from video."""
        if not MOVIEPY_AVAILABLE:
            raise ImportError("moviepy is not installed")
        
        clip = VideoFileClip(video_path)
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        clip.audio.write_audiofile(str(output_file))
        clip.close()
        
        return str(output_file)
    
    # ==================== CONVERSION UTILITIES ====================
    
    def convert_document(self, input_path: str, output_path: str, 
                         input_format: str = None, output_format: str = None) -> str:
        """
        Convert documents between formats.
        Supports: text, PDF, HTML, Markdown, images, etc.
        """
        input_path = Path(input_path)
        output_path = Path(output_path)
        
        # Auto-detect formats if not specified
        if not input_format:
            input_format = input_path.suffix.lower()
        if not output_format:
            output_format = output_path.suffix.lower()
        
        output_file = self.output_dir / output_path
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Text to PDF
        if input_format in ['.txt', '.text'] and output_format == '.pdf':
            content = self.read_text_file(str(input_path))
            return self.create_pdf(str(output_file), 
                                  [{'type': 'paragraph', 'text': content}])
        
        # Markdown to PDF via HTML
        elif input_format == '.md' and output_format == '.pdf':
            content = self.read_text_file(str(input_path))
            html_content = markdown.markdown(content)
            temp_html = self.temp_dir / 'temp.html'
            with open(temp_html, 'w') as f:
                f.write(html_content)
            # Use external tool like wkhtmltopdf if available
            try:
                subprocess.run(['wkhtmltopdf', str(temp_html), str(output_file)], check=True)
            except:
                # Fallback to creating a PDF with reportlab
                return self.create_pdf(str(output_file), 
                                      [{'type': 'paragraph', 'text': content}])
        
        # Word to PDF
        elif input_format in ['.docx', '.doc'] and output_format == '.pdf':
            # This requires additional tools, but we'll implement basic version
            if DOCX_AVAILABLE:
                doc_content = self.read_word_document(str(input_path))
                combined_text = '\n\n'.join(doc_content['paragraphs'])
                return self.create_pdf(str(output_file),
                                      [{'type': 'paragraph', 'text': combined_text}])
        
        # Default: copy with conversion using pandas for spreadsheets
        elif input_format in ['.xlsx', '.csv', '.xls'] and output_format in ['.csv', '.xlsx']:
            df = pd.read_excel(str(input_path)) if input_format != '.csv' else pd.read_csv(str(input_path))
            if output_format == '.csv':
                df.to_csv(str(output_file), index=False)
            else:
                df.to_excel(str(output_file), index=False)
            return str(output_file)
        
        # Image conversion
        elif input_format in ['.jpg', '.jpeg', '.png', '.gif', '.bmp'] and output_format in ['.jpg', '.png', '.webp']:
            return self.convert_image(str(input_path), str(output_file), output_format[1:])
        
        else:
            raise NotImplementedError(f"Conversion from {input_format} to {output_format} not implemented")
        
        return str(output_file)
    
    # ==================== BATCH PROCESSING ====================
    
    def batch_process(self, input_dir: str, output_dir: str, 
                     operation: str, **kwargs) -> List[str]:
        """
        Batch process multiple files in a directory.
        
        Args:
            input_dir: Directory containing input files
            output_dir: Directory for output files
            operation: Operation to perform ('convert', 'extract_text', 'resize', etc.)
            **kwargs: Additional parameters for the operation
        """
        input_path = Path(input_dir)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        results = []
        
        for file_path in input_path.iterdir():
            if file_path.is_file():
                try:
                    if operation == 'convert':
                        output_format = kwargs.get('output_format', '.txt')
                        output_file = output_path / f"{file_path.stem}{output_format}"
                        result = self.convert_document(str(file_path), str(output_file))
                        results.append(result)
                    
                    elif operation == 'extract_text':
                        text = self.extract_text_from_file(str(file_path))
                        output_file = output_path / f"{file_path.stem}.txt"
                        with open(output_file, 'w', encoding='utf-8') as f:
                            f.write(text)
                        results.append(str(output_file))
                    
                    elif operation == 'resize_images' and file_path.suffix.lower() in ['.jpg', '.png']:
                        output_file = output_path / file_path.name
                        result = self.resize_image(str(file_path), str(output_file), 
                                                  width=kwargs.get('width'), 
                                                  height=kwargs.get('height'))
                        results.append(result)
                
                except Exception as e:
                    print(f"Error processing {file_path}: {e}")
        
        return results
    
    def extract_text_from_file(self, file_path: str) -> str:
        """Universal text extraction from any supported document type."""
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        if suffix in ['.txt', '.text', '.log', '.asc']:
            return self.read_text_file(str(file_path))
        
        elif suffix in ['.docx', '.doc']:
            if DOCX_AVAILABLE:
                result = self.read_word_document(str(file_path))
                return '\n'.join(result['paragraphs'])
        
        elif suffix == '.pdf':
            result = self.read_pdf(str(file_path), extract_tables=False)
            return '\n'.join([page['text'] for page in result['pages']])
        
        elif suffix in ['.xlsx', '.xls', '.csv', '.tsv']:
            df = pd.read_excel(str(file_path)) if suffix != '.csv' else pd.read_csv(str(file_path))
            return df.to_string()
        
        elif suffix in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            if EASYOCR_AVAILABLE:
                return self.extract_text_from_image(str(file_path))
            else:
                return "[OCR not available. Install EasyOCR: pip install easyocr]"
        
        elif suffix in ['.html', '.htm']:
            if BEAUTIFULSOUP_AVAILABLE:
                with open(file_path, 'r', encoding='utf-8') as f:
                    soup = BeautifulSoup(f.read(), 'lxml')
                    return soup.get_text()
        
        elif suffix == '.json':
            data = self.read_json(str(file_path))
            return json.dumps(data, indent=2)
        
        else:
            return f"Unsupported format: {suffix}"
    
    def cleanup_temp(self):
        """Remove temporary files."""
        shutil.rmtree(self.temp_dir, ignore_errors=True)
    
    def __del__(self):
        """Cleanup on destruction."""
        self.cleanup_temp()


# ==================== CONVENIENCE FUNCTIONS ====================

def create_document_processor(output_dir: str = "processed_docs") -> DocumentProcessor:
    """Factory function to create a DocumentProcessor instance."""
    return DocumentProcessor(output_dir)


# ==================== USAGE EXAMPLES ====================

if __name__ == "__main__":
    # Example usage
    processor = create_document_processor("example_output")
    
    print("Document Processor initialized. Available formats:", 
          list(processor.supported_formats.keys()))
    
    # Example: Create a text file
    processor.create_text_file("Hello, World!\nThis is a test.", "test.txt")
    print("Created test.txt")
    
    # Example: Create a JSON file
    data = {"name": "Test", "version": 1.0, "data": [1, 2, 3, 4, 5]}
    processor.write_json(data, "test.json")
    print("Created test.json")
    
    # Example: Read back the JSON
    read_data = processor.read_json("example_output/test.json")
    print("Read JSON:", read_data)
    
    # Example: Create a CSV file
    csv_data = [
        {"name": "Alice", "age": 30, "city": "New York"},
        {"name": "Bob", "age": 25, "city": "Los Angeles"},
        {"name": "Charlie", "age": 35, "city": "Chicago"}
    ]
    processor.create_csv(csv_data, "test.csv")
    print("Created test.csv")
    
    print("\nAll examples completed successfully!")